"""Build conference presentation from the ICIE template."""
import copy
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

TEMPLATE = "slides/Oral presentation ICIE.pptx"
OUT = "slides/presentation.pptx"
FIG = "paper/docx_figures"

SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)


def clone_slide(prs, src_slide):
    """Duplicate a slide (deep copy of XML + relationships)."""
    src_part = src_slide.part
    new_slide = prs.slides.add_slide(src_slide.slide_layout)
    # Remove placeholders auto-added by layout
    for sh in list(new_slide.shapes):
        sh._element.getparent().remove(sh._element)
    # Copy spTree children from src
    src_spTree = src_part._element.find(qn("p:cSld")).find(qn("p:spTree"))
    new_spTree = new_slide.part._element.find(qn("p:cSld")).find(qn("p:spTree"))
    for child in list(src_spTree):
        tag = etree.QName(child.tag).localname
        if tag in ("nvGrpSpPr", "grpSpPr"):
            continue
        new_spTree.append(copy.deepcopy(child))
    # Copy relationships (for the picture)
    for rel in src_part.rels.values():
        if rel.reltype.endswith("/image"):
            new_slide.part.relate_to(rel.target_part, rel.reltype)
    return new_slide


def set_run_text(tf, paragraphs):
    """Replace text frame content. paragraphs: list of (text, dict opts)."""
    # Clear all but first paragraph
    p_elems = tf._txBody.findall(qn("a:p"))
    for p in p_elems[1:]:
        tf._txBody.remove(p)
    # Clear runs in first paragraph
    first_p = p_elems[0]
    for r in first_p.findall(qn("a:r")):
        first_p.remove(r)
    for br in first_p.findall(qn("a:br")):
        first_p.remove(br)
    # Helper to add a paragraph with bullet text
    from pptx.util import Pt as _Pt
    for i, (text, opts) in enumerate(paragraphs):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.text = text
        if "size" in opts:
            for r in para.runs:
                r.font.size = _Pt(opts["size"])
        if opts.get("bold"):
            for r in para.runs:
                r.font.bold = True
        if "color" in opts:
            for r in para.runs:
                r.font.color.rgb = RGBColor(*opts["color"])
        if "align" in opts:
            from pptx.enum.text import PP_ALIGN
            para.alignment = opts["align"]


def add_textbox(slide, left, top, width, height, paragraphs, font_name="Calibri"):
    from pptx.enum.text import MSO_AUTO_SIZE
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    for i, (text, opts) in enumerate(paragraphs):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        run = para.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(opts.get("size", 18))
        if opts.get("bold"):
            run.font.bold = True
        if "color" in opts:
            run.font.color.rgb = RGBColor(*opts["color"])
        if "align" in opts:
            from pptx.enum.text import PP_ALIGN
            para.alignment = opts["align"]
        if opts.get("bullet"):
            # add bullet via XML
            pPr = para._pPr if para._pPr is not None else para._p.get_or_add_pPr()
            buChar = etree.SubElement(pPr, qn("a:buChar"))
            buChar.set("char", "•")
    return tb


def set_title(slide, ru, en):
    """Set the title shape on a cloned slide (Прямоугольник 3)."""
    for sh in slide.shapes:
        if sh.has_text_frame and "Заголовок слайда" in sh.text_frame.text or (
            sh.name.startswith("Прямоугольник") and sh.has_text_frame and sh.top < Inches(1)
        ):
            tf = sh.text_frame
            # keep the existing styling; replace text in two paragraphs
            p_elems = tf._txBody.findall(qn("a:p"))
            # Clear extras
            for p in p_elems[2:]:
                tf._txBody.remove(p)
            # First paragraph
            p0 = p_elems[0]
            for r in p0.findall(qn("a:r")):
                p0.remove(r)
            for br in p0.findall(qn("a:br")):
                p0.remove(br)
            r = etree.SubElement(p0, qn("a:r"))
            rPr = etree.SubElement(r, qn("a:rPr"))
            rPr.set("lang", "ru-RU")
            rPr.set("sz", "1800")
            rPr.set("b", "1")
            rPr.set("dirty", "0")
            t = etree.SubElement(r, qn("a:t"))
            t.text = ru + " / " + en
            return


def remove_textbox_138(slide):
    """Remove the 'Уважаемые участники' instructional text box."""
    for sh in list(slide.shapes):
        if sh.has_text_frame and "Уважаемые участники" in sh.text_frame.text:
            sh._element.getparent().remove(sh._element)
            return


def add_picture_centered(slide, img_path, top, max_w, max_h):
    from PIL import Image
    im = Image.open(img_path)
    iw, ih = im.size
    ratio = iw / ih
    if max_w / max_h > ratio:
        h = max_h
        w = int(max_h * ratio)
    else:
        w = max_w
        h = int(max_w / ratio)
    left = int((SLIDE_W - w) / 2)
    slide.shapes.add_picture(img_path, left, top, width=w, height=h)


def main():
    prs = Presentation(TEMPLATE)
    slides = list(prs.slides)
    cover = slides[0]
    template_content = slides[1]  # has background pic + title + instructional textbox

    # === Slide 1: Cover ===
    # Shape [1] = title rectangle ('Прямоугольник 13'), Shape [2] = name rectangle
    for sh in cover.shapes:
        if sh.has_text_frame:
            txt = sh.text_frame.text
            if "НАЗВАНИЕ ДОКЛАДА" in txt:
                tf = sh.text_frame
                p_elems = tf._txBody.findall(qn("a:p"))
                for p in p_elems[2:]:
                    tf._txBody.remove(p)
                # Russian title
                p0 = p_elems[0]
                for r in p0.findall(qn("a:r")):
                    p0.remove(r)
                for br in p0.findall(qn("a:br")):
                    p0.remove(br)
                r = etree.SubElement(p0, qn("a:r"))
                rPr = etree.SubElement(r, qn("a:rPr"))
                rPr.set("lang", "ru-RU"); rPr.set("sz", "2400"); rPr.set("b", "1"); rPr.set("dirty", "0")
                t = etree.SubElement(r, qn("a:t"))
                t.text = "Сравнительное исследование классических, машинно-обучаемых, глубоких и фундаментальных моделей для краткосрочного прогноза электропотребления"
                # English title
                p1 = p_elems[1]
                for r in p1.findall(qn("a:r")):
                    p1.remove(r)
                for br in p1.findall(qn("a:br")):
                    p1.remove(br)
                r = etree.SubElement(p1, qn("a:r"))
                rPr = etree.SubElement(r, qn("a:rPr"))
                rPr.set("lang", "en-US"); rPr.set("sz", "2000"); rPr.set("b", "1"); rPr.set("dirty", "0")
                t = etree.SubElement(r, qn("a:t"))
                t.text = "A Comparative Study of Machine Learning, Deep Learning, and Foundation Models for Short-Term Electricity Load Forecasting"
            elif "ФИО докладчика" in txt:
                tf = sh.text_frame
                p_elems = tf._txBody.findall(qn("a:p"))
                for p in p_elems[2:]:
                    tf._txBody.remove(p)
                p0 = p_elems[0]
                for r in p0.findall(qn("a:r")):
                    p0.remove(r)
                for br in p0.findall(qn("a:br")):
                    p0.remove(br)
                r = etree.SubElement(p0, qn("a:r"))
                rPr = etree.SubElement(r, qn("a:rPr"))
                rPr.set("lang", "ru-RU"); rPr.set("sz", "1600"); rPr.set("dirty", "0")
                t = etree.SubElement(r, qn("a:t"))
                t.text = "Князев А. М., Гречаников Н. Д."
                p1 = p_elems[1]
                for r in p1.findall(qn("a:r")):
                    p1.remove(r)
                for br in p1.findall(qn("a:br")):
                    p1.remove(br)
                r = etree.SubElement(p1, qn("a:r"))
                rPr = etree.SubElement(r, qn("a:rPr"))
                rPr.set("lang", "en-US"); rPr.set("sz", "1600"); rPr.set("dirty", "0")
                t = etree.SubElement(r, qn("a:t"))
                t.text = "A. M. Kniazev, N. D. Grechanikov"
                # Add affiliation as new paragraph
                p2 = etree.SubElement(tf._txBody, qn("a:p"))
                r = etree.SubElement(p2, qn("a:r"))
                rPr = etree.SubElement(r, qn("a:rPr"))
                rPr.set("lang", "en-US"); rPr.set("sz", "1400"); rPr.set("i", "1"); rPr.set("dirty", "0")
                t = etree.SubElement(r, qn("a:t"))
                t.text = "RUDN University, Moscow, Russia"

    # Clone the pristine template slide 7 more times BEFORE editing any of them
    extra = [clone_slide(prs, template_content) for _ in range(7)]

    # === Slide 2: Motivation ===
    s2 = template_content
    set_title(s2, "Мотивация", "Motivation")
    remove_textbox_138(s2)
    add_textbox(s2, Inches(0.6), Inches(1.4), Inches(8.8), Inches(5.6), [
        ("Краткосрочный прогноз потребления электроэнергии — базовая задача планирования работы энергосистемы: выбор состава генерирующего оборудования, расчёт резерва мощности, управление спросом.",
         {"size": 18, "bullet": True}),
        ("", {"size": 8}),
        ("За последний год появилось новое поколение фундаментальных моделей для временных рядов — Chronos, TimesFM, Lag-Llama. Они предобучены на больших корпусах и дают прогноз без дообучения на целевом ряде.",
         {"size": 18, "bullet": True}),
        ("", {"size": 8}),
        ("Однако в литературе почти нет прямых сравнений фундаментальных моделей со специально обученными нейросетевыми и классическими методами по единому протоколу без утечек данных.",
         {"size": 18, "bullet": True}),
        ("", {"size": 10}),
        ("Вопрос: насколько прогноз без дообучения конкурентоспособен против локально обученной модели?",
         {"size": 20, "bold": True, "color": (0xC0, 0x50, 0x4D)}),
    ])

    # === Slide 3: Aim & Tasks ===
    s3 = extra[0]
    set_title(s3, "Цель и задачи", "Aim and Tasks")
    remove_textbox_138(s3)
    add_textbox(s3, Inches(0.6), Inches(1.4), Inches(8.8), Inches(5.6), [
        ("Цель", {"size": 24, "bold": True, "color": (0x1F, 0x49, 0x7D)}),
        ("Количественно оценить, как фундаментальные модели без дообучения сравниваются со специально обученными нейросетевыми и классическими методами по единому хронологическому протоколу без утечек данных.",
         {"size": 18}),
        ("", {"size": 10}),
        ("Задачи", {"size": 24, "bold": True, "color": (0x1F, 0x49, 0x7D)}),
        ("1. Построить единый конвейер обработки: хронологическое разбиение, стандартизация по обучающей выборке, контекст 336 ч, горизонты 24 / 96 / 168 ч.",
         {"size": 18}),
        ("2. Сравнить 8 моделей из 3 семейств на двух наборах данных (ECL и ETTh1/HUFL) по 5 случайным инициализациям.",
         {"size": 18}),
        ("3. Провести абляционное исследование по бюджету обучения (от 500 до 5000 шагов).",
         {"size": 18}),
        ("4. Проанализировать компромисс между точностью прогноза и задержкой инференса.",
         {"size": 18}),
    ])

    # === Slide 4: Protocol & Data ===
    s4 = extra[1]
    set_title(s4, "Протокол и данные", "Protocol and Data")
    remove_textbox_138(s4)
    add_textbox(s4, Inches(0.4), Inches(1.4), Inches(4.6), Inches(5.6), [
        ("Данные", {"size": 20, "bold": True, "color": (0x1F, 0x49, 0x7D)}),
        ("• ECL — суммарная нагрузка сети", {"size": 16}),
        ("• ETTh1 / HUFL — сигнал датчика трансформаторной подстанции", {"size": 16}),
        ("", {"size": 6}),
        ("Разбиение выборки", {"size": 20, "bold": True, "color": (0x1F, 0x49, 0x7D)}),
        ("• 70 / 15 / 15 в хронологическом порядке", {"size": 16}),
        ("• Стандартизация по обучающей части", {"size": 16}),
        ("• Скользящие окна без утечек данных", {"size": 16}),
        ("", {"size": 6}),
        ("Параметры эксперимента", {"size": 20, "bold": True, "color": (0x1F, 0x49, 0x7D)}),
        ("• Контекст L = 336 ч (две недели)", {"size": 16}),
        ("• Горизонты H ∈ {24, 96, 168} ч", {"size": 16}),
        ("• 5 случайных инициализаций на конфигурацию", {"size": 16}),
    ])
    add_picture_centered(s4, f"{FIG}/fig1_methodology.png", Inches(1.0), Inches(4.8), Inches(6.2))
    # Move pic to right side
    pic = s4.shapes[-1]
    pic.left = Inches(5.1)
    pic.top = Inches(1.1)
    pic.width = Inches(4.7)
    pic.height = Inches(4.4)

    # === Slide 5: Headline ECL result ===
    s5 = extra[2]
    set_title(s5, "Результаты на ECL", "Results on ECL")
    remove_textbox_138(s5)
    add_textbox(s5, Inches(0.4), Inches(1.3), Inches(9.2), Inches(0.6), [
        ("Обученные нейросети: MAE ≈ 26–32 тыс.  ·  фундаментальные модели без дообучения: MAE ≈ 260 тыс. — разрыв почти на порядок.",
         {"size": 14, "bold": True, "color": (0xC0, 0x50, 0x4D)}),
    ])
    add_picture_centered(s5, f"{FIG}/fig6_ecl_accuracy.png", Inches(2.0), Inches(8.4), Inches(4.3))
    add_textbox(s5, Inches(0.4), Inches(6.4), Inches(9.2), Inches(1.0), [
        ("ECL, H = 24:  iTransformer 26 083 · DLinear 27 038 · PatchTST 31 430 · XGBoost 31 738 · SARIMA 263 914 · TimesFM 259 345 · Chronos-Bolt 259 546.",
         {"size": 11}),
        ("iTransformer лучший на ECL на всех горизонтах; DLinear — стабильно второй на длинных.",
         {"size": 11}),
    ])

    # === Slide 6: ETTh1 — flip ===
    s6 = extra[3]
    set_title(s6, "Результаты на ETTh1 / HUFL", "Results on ETTh1 / HUFL")
    remove_textbox_138(s6)
    add_textbox(s6, Inches(0.4), Inches(1.3), Inches(9.2), Inches(0.6), [
        ("На ETTh1, горизонт 24 ч, лучший — Chronos-Bolt (MAE 3.00). Фундаментальная модель не просто конкурентна, а лидер.",
         {"size": 14, "bold": True, "color": (0x1F, 0x49, 0x7D)}),
    ])
    add_picture_centered(s6, f"{FIG}/fig7_etth1_accuracy.png", Inches(2.0), Inches(8.4), Inches(4.3))
    add_textbox(s6, Inches(0.4), Inches(6.4), Inches(9.2), Inches(1.0), [
        ("ETTh1, H = 24:  Chronos-Bolt 3.00 · DLinear 3.07 · PatchTST 3.10 · TimesFM 3.16.",
         {"size": 11}),
        ("На H = 96 и H = 168 вперёд выходит DLinear, фундаментальные модели — в пределах нескольких процентов.",
         {"size": 11}),
        ("Вывод: прогноз без дообучения работает там, где целевой ряд близок по характеру к данным предобучения.",
         {"size": 11, "bold": True}),
    ])

    # === Slide 7: Training-budget ablation ===
    s8 = extra[4]
    set_title(s8, "Абляция по бюджету обучения", "Training-budget ablation")
    remove_textbox_138(s8)
    add_textbox(s8, Inches(0.4), Inches(1.3), Inches(9.2), Inches(0.6), [
        ("Разрыв на ECL — не следствие недообучения нейросетей, а собственное ограничение режима без дообучения.",
         {"size": 14, "bold": True, "color": (0xC0, 0x50, 0x4D)}),
    ])
    add_picture_centered(s8, f"{FIG}/fig5_horizon.png", Inches(2.0), Inches(8.4), Inches(4.3))
    add_textbox(s8, Inches(0.4), Inches(6.4), Inches(9.2), Inches(1.0), [
        ("ECL, H = 24: PatchTST и iTransformer обучались по {500, 1000, 2000, 5000} шагам с ранней остановкой.",
         {"size": 11}),
        ("PatchTST 31.4 – 33.7 тыс. MAE · iTransformer 26.1 – 29.8 тыс. MAE · разброс между бюджетами менее 4 %.",
         {"size": 11}),
        ("MAE фундаментальных моделей остаётся около 260 тыс. независимо от бюджета обучения нейросетей.",
         {"size": 11}),
    ])

    # === Slide 8: Accuracy-latency Pareto ===
    s9 = extra[5]
    set_title(s9, "Точность и задержка инференса", "Accuracy and inference latency")
    remove_textbox_138(s9)
    add_textbox(s9, Inches(0.4), Inches(1.3), Inches(9.2), Inches(0.6), [
        ("Нейросети лидируют по точности, классические методы — по скорости; фундаментальные модели на ECL оказались в неудачной «середине».",
         {"size": 13}),
    ])
    add_picture_centered(s9, f"{FIG}/fig4_pareto.png", Inches(2.0), Inches(8.4), Inches(4.3))
    add_textbox(s9, Inches(0.4), Inches(6.4), Inches(9.2), Inches(1.0), [
        ("SARIMA — менее 6 мс на окно, самая быстрая точка фронта Парето.",
         {"size": 11}),
        ("Фундаментальные модели в 13.6 раза медленнее SARIMA и в 4.2 раза тяжелее DLinear по числу параметров.",
         {"size": 11}),
        ("Рекомендуемое многоуровневое развёртывание: SARIMA — при жёстком ограничении задержки; нейросети — для основной нагрузки; фундаментальные модели — для холодного старта.",
         {"size": 11, "bold": True}),
    ])

    # === Slide 9: Conclusion ===
    s10 = extra[6]
    set_title(s10, "Выводы", "Conclusion")
    remove_textbox_138(s10)
    add_textbox(s10, Inches(0.6), Inches(1.4), Inches(8.8), Inches(5.6), [
        ("1. По единому протоколу без утечек данных обученные нейросети (iTransformer, DLinear, PatchTST) доминируют на ECL; XGBoost — сильнейший метод без глубокого обучения.",
         {"size": 17}),
        ("", {"size": 6}),
        ("2. Фундаментальные модели без дообучения отстают на ECL примерно на порядок, но конкурентны и даже лучшие на ETTh1 при горизонте 24 ч.",
         {"size": 17}),
        ("", {"size": 6}),
        ("3. Разрыв на ECL — не следствие недостаточного обучения нейросетей, а собственное ограничение режима без дообучения при существенном сдвиге распределения.",
         {"size": 17}),
        ("", {"size": 6}),
        ("4. Исходный код, конфигурации и все 144 строки результатов опубликованы для воспроизведения.",
         {"size": 17}),
        ("", {"size": 10}),
        ("Спасибо за внимание!",
         {"size": 22, "bold": True, "align": __import__("pptx").enum.text.PP_ALIGN.CENTER, "color": (0x1F, 0x49, 0x7D)}),
    ])

    prs.save(OUT)
    print(f"Saved: {OUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
